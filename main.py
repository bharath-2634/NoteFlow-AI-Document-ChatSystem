from fastapi import FastAPI, UploadFile, File, HTTPException
from dotenv import load_dotenv
from sentence_transformers import SentenceTransformer
import os
from uuid import uuid4
import io

from PyPDF2 import PdfReader
import docx
from pptx import Presentation

from fastapi import Query
from pydantic import BaseModel

from langchain.text_splitter import RecursiveCharacterTextSplitter

from pinecone import Pinecone, PodSpec

load_dotenv()
app = FastAPI()


model = SentenceTransformer('all-MiniLM-L6-v2')

pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
index_name = os.getenv("PINECONE_INDEX_NAME")

index = pc.Index(index_name)

def extract_text_from_pptx(file_bytes: io.BytesIO) -> str:
    from pptx import Presentation
    prs = Presentation(file_bytes)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_text(file: UploadFile) -> str:
    ext = os.path.splitext(file.filename)[1].lower()
    file_bytes = io.BytesIO(file.file.read())

    if ext == ".pdf":
        reader = PdfReader(file_bytes)
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    elif ext == ".docx":
        doc = docx.Document(file_bytes)
        return "\n".join([p.text for p in doc.paragraphs])
    elif ext == ".pptx":
        return extract_text_from_pptx(file_bytes)
    elif ext == ".txt":
        return file_bytes.read().decode("utf-8")
    else:
        raise HTTPException(status_code=400, detail="Unsupported file format")

def chunk_text(text: str):
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    return splitter.split_text(text)

def embed_chunks(chunks):
    return model.encode(chunks).tolist()

@app.post("/embed")
async def embed_file(file: UploadFile = File(...)):
    try:
        text = extract_text(file)
        chunks = chunk_text(text)
        embeddings = embed_chunks(chunks)

        vectors = [{
            "id": str(uuid4()),
            "values": emb,
            "metadata": {"text": chunk}
        } for chunk, emb in zip(chunks, embeddings)]

        index.upsert(vectors=vectors)

        return {
            "success": True,
            "chunks_indexed": len(chunks)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


class QueryRequest(BaseModel):
    query: str

@app.post("/query")
async def query_docs(request: QueryRequest):
    query_embedding = model.encode(request.query).tolist()

    # Search Pinecone
    response = index.query(
        vector=query_embedding,
        top_k=5,
        include_metadata=True
    )

    matches = response.get("matches", [])

    results = [
        {
            "score": match["score"],
            "text": match["metadata"]["text"]
        }
        for match in matches
    ]

    return {
        "success": True,
        "results": results
    }
